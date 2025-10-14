from __future__ import annotations
from pathlib import Path
from typing import List, Dict
import json, re
from io import BytesIO

import fitz  # PyMuPDF

import os
import requests
from pptx import Presentation
from docx import Document  # Add this import at the top
# --- helpers ---
def _approx_tokens(s) -> int:
    return max(1, len(s) // 4)  # rough heuristic

def _split_into_chunks(text: str, max_chars: int = 1500) -> List[str]:
    paras = re.split(r"\n\s*\n", (text or "").strip())
    out, buf = [], ""
    for p in paras:
        p = p.strip()
        if not p:
            continue
        if len(buf) + len(p) + 2 <= max_chars:
            buf = (buf + "\n\n" + p).strip() if buf else p
        else:
            if buf:
                out.append(buf)
            buf = p
    if buf:
        out.append(buf)
    return out

# --- per-type parsers ---
def _parse_pdf(path, name) -> List[Dict]:
    doc = fitz.open(path)
    chunks = []
    for i, page in enumerate(doc, start=1):
        text = page.get_text("text") or ""
        for j, c in enumerate(_split_into_chunks(text), start=1):
            chunks.append({
                "id": f"{name}::p{i}::{j}",
                "text": c,
                "tokens": _approx_tokens(c),
                "doc_type": "pdf",
                "tags": [],
                "source": {"file": name, "locator": f"p{i}"}
            })
    doc.close()
    return chunks

def _parse_docx(filepath, name) -> List[Dict]:
    """filepath: string path to DOCX file"""
    doc = Document(filepath)
    
    # Extract all text from paragraphs
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
    text = "\n\n".join(paragraphs)
    
    # Also extract text from tables if any
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                text += "\n" + row_text
    
    
    chunks = []
    for j, c in enumerate(_split_into_chunks(text), start=1):
        chunks.append({
            "id": f"{name}::doc::{j}",
            "text": c,
            "tokens": _approx_tokens(c),
            "doc_type": "docx",
            "tags": [],
            "source": {"file": name, "locator": f"sec{j}"}
        })
    return chunks
        

def _parse_pptx(path, name) -> List[Dict]:
    prs = Presentation(path)
    chunks = []
    def text_from_shape(shape) -> str:
        if hasattr(shape, "text"):
            return shape.text or ""
        if hasattr(shape, "text_frame") and shape.text_frame:
            return "\n".join(p.text or "" for p in shape.text_frame.paragraphs)
        return ""
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            t = text_from_shape(shape).strip()
            if t:
                texts.append(t)
        slide_text = "\n".join(texts).strip()
        if not slide_text:
            continue
        for j, c in enumerate(_split_into_chunks(slide_text, max_chars=1200), start=1):
            chunks.append({
                "id": f"{name}::s{i}::{j}",
                "text": c,
                "tokens": _approx_tokens(c),
                "doc_type": "pptx",
                "tags": [],
                "source": {"file": name, "locator": f"s{i}"}
            })
    return chunks

def _parse_txt(path, name) -> List[Dict]:
    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
        text = f.read()
    chunks = []
    for j, c in enumerate(_split_into_chunks(text), start=1):
        chunks.append({
            "id": f"{name}::txt::{j}",
            "text": c,
            "tokens": _approx_tokens(c),
            "doc_type": "txt",
            "tags": [],
            "source": {"file": name, "locator": f"sec{j}"}
        })
    return chunks

PARSERS = {
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".pptx": _parse_pptx,
    ".txt": _parse_txt,
}

# In parsing.py - ingest_files function
def ingest_files(filenames: List[Dict]) -> bytes:
    all_chunks: List[Dict] = []
    
    for file_info in filenames:
        filename = file_info['filename']
        filepath = file_info['path']
        ext = "." + filename.split(".")[-1].lower()
        
        # DEBUG: Only check DOCX files
        
        # Read from local file
        if ext == ".docx":
            import time, zipfile

            time.sleep(0.05)  # small wait to avoid race (Lambda-specific quirk)
            print("try zip")
            with open(filepath, "rb") as f:
                # double-check ZIP validity by actually opening it as zip
                z = zipfile.ZipFile(f)
                z.testzip()  # will raise BadZipFile if something is wrong
            print("try docx")
            # Now open with python-docx
            doc = Document(filepath)
            print("success?")
        parser = PARSERS.get(ext)
        if parser:
            all_chunks.extend(parser(filepath, filename))
    
    return all_chunks
